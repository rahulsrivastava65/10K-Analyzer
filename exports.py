from __future__ import annotations

from io import BytesIO
from typing import Any, Dict, List, Sequence

import pandas as pd
from openpyxl.chart import BarChart, LineChart, Reference
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.shapes import Drawing, String
from reportlab.graphics.widgets.markers import makeMarker
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


def _fmt_currency(value: float | None) -> str:
    if value is None or pd.isna(value):
        return "-"
    negative = value < 0
    amount = abs(float(value))
    if amount >= 1_000_000_000:
        text = f"${amount / 1_000_000_000:,.2f}B"
    elif amount >= 1_000_000:
        text = f"${amount / 1_000_000:,.2f}M"
    else:
        text = f"${amount:,.0f}"
    return f"({text})" if negative else text


def _fmt_percent(value: float | None) -> str:
    if value is None or pd.isna(value):
        return "-"
    return f"{float(value):.1f}%"


def _fmt_ratio(value: float | None) -> str:
    if value is None or pd.isna(value):
        return "-"
    return f"{float(value):.2f}x"


def _history_years(history: pd.DataFrame) -> List[str]:
    if history.empty or "Year" not in history.columns:
        return []
    return [f"FY{int(year)}" for year in history["Year"].tolist()]


def _history_billions(history: pd.DataFrame, metric: str) -> List[float]:
    if history.empty or metric not in history.columns:
        return []
    return [0.0 if pd.isna(value) else float(value) / 1_000_000_000 for value in history[metric].tolist()]


def _add_excel_charts(writer: pd.ExcelWriter, history: pd.DataFrame) -> None:
    if history.empty or "Year" not in history.columns:
        return

    workbook = writer.book
    history_sheet = writer.sheets["History"]
    charts_sheet = workbook.create_sheet("Charts")
    charts_sheet["B1"] = "10K Summary Charts"

    max_row = len(history) + 1
    year_col = history.columns.get_loc("Year") + 1

    def add_line_chart(title: str, metrics: List[str], anchor: str) -> None:
        available = [metric for metric in metrics if metric in history.columns]
        if not available:
            return
        chart = LineChart()
        chart.title = title
        chart.style = 13
        chart.height = 8
        chart.width = 14
        chart.y_axis.title = "$B"
        chart.x_axis.title = "Fiscal Year"
        data = Reference(
            history_sheet,
            min_col=history.columns.get_loc(available[0]) + 1,
            max_col=history.columns.get_loc(available[-1]) + 1,
            min_row=1,
            max_row=max_row,
        )
        cats = Reference(history_sheet, min_col=year_col, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.legend.position = "r"
        charts_sheet.add_chart(chart, anchor)

    def add_bar_chart(title: str, metrics: List[str], anchor: str) -> None:
        available = [metric for metric in metrics if metric in history.columns]
        if not available:
            return
        chart = BarChart()
        chart.type = "col"
        chart.grouping = "clustered"
        chart.overlap = 0
        chart.title = title
        chart.style = 10
        chart.height = 8
        chart.width = 14
        chart.y_axis.title = "$B"
        chart.x_axis.title = "Fiscal Year"
        data = Reference(
            history_sheet,
            min_col=history.columns.get_loc(available[0]) + 1,
            max_col=history.columns.get_loc(available[-1]) + 1,
            min_row=1,
            max_row=max_row,
        )
        cats = Reference(history_sheet, min_col=year_col, min_row=2, max_row=max_row)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        chart.legend.position = "r"
        charts_sheet.add_chart(chart, anchor)

    add_line_chart("Revenue & Net Income Trend", ["Revenue", "Net Income"], "B3")
    add_line_chart("Expenses Trend", ["Expenses"], "B20")
    add_bar_chart("Assets vs Liabilities", ["Assets", "Liabilities"], "J3")


def _build_pdf_line_chart(
    history: pd.DataFrame,
    title: str,
    metrics: Sequence[str],
    series_colors: Sequence[str],
) -> Drawing | None:
    years = _history_years(history)
    series = [_history_billions(history, metric) for metric in metrics]
    if len(years) < 2 or any(not values for values in series):
        return None
    all_values = [value for values in series for value in values]

    drawing = Drawing(520, 210)
    drawing.add(String(12, 190, f"{title} ($B)", fontName="Helvetica-Bold", fontSize=11, fillColor=colors.HexColor("#17365d")))

    chart = HorizontalLineChart()
    chart.x = 42
    chart.y = 40
    chart.height = 120
    chart.width = 440
    chart.data = [tuple(values) for values in series]
    chart.categoryAxis.categoryNames = years
    chart.categoryAxis.labels.boxAnchor = "n"
    chart.categoryAxis.labels.fontName = "Helvetica"
    chart.categoryAxis.labels.fontSize = 8
    chart.valueAxis.labels.fontName = "Helvetica"
    chart.valueAxis.labels.fontSize = 8
    chart.valueAxis.valueMin = min(0, min(all_values))
    chart.valueAxis.valueMax = max(all_values) * 1.15 if max(all_values) else 1
    chart.valueAxis.valueStep = max(chart.valueAxis.valueMax / 4, 1)
    for idx, color in enumerate(series_colors):
        chart.lines[idx].strokeColor = colors.HexColor(color)
        chart.lines[idx].strokeWidth = 2
        chart.lines[idx].symbol = makeMarker("FilledCircle")
    drawing.add(chart)
    return drawing


def _build_pdf_bar_chart(history: pd.DataFrame, title: str, left_metric: str, right_metric: str) -> Drawing | None:
    years = _history_years(history)
    left_values = _history_billions(history, left_metric)
    right_values = _history_billions(history, right_metric)
    if len(years) < 2 or not left_values or not right_values:
        return None

    drawing = Drawing(520, 210)
    drawing.add(String(12, 190, f"{title} ($B)", fontName="Helvetica-Bold", fontSize=11, fillColor=colors.HexColor("#17365d")))

    chart = VerticalBarChart()
    chart.x = 45
    chart.y = 40
    chart.height = 120
    chart.width = 430
    chart.data = [tuple(left_values), tuple(right_values)]
    chart.categoryAxis.categoryNames = years
    chart.categoryAxis.labels.boxAnchor = "n"
    chart.categoryAxis.labels.fontName = "Helvetica"
    chart.categoryAxis.labels.fontSize = 8
    chart.valueAxis.labels.fontName = "Helvetica"
    chart.valueAxis.labels.fontSize = 8
    chart.valueAxis.valueMin = 0 if min(left_values + right_values) >= 0 else min(left_values + right_values) * 1.1
    chart.valueAxis.valueMax = max(left_values + right_values) * 1.15 if max(left_values + right_values) else 1
    chart.valueAxis.valueStep = max(chart.valueAxis.valueMax / 4, 1)
    chart.barWidth = 8
    chart.groupSpacing = 10
    chart.bars[0].fillColor = colors.HexColor("#17365d")
    chart.bars[1].fillColor = colors.HexColor("#7a8795")
    drawing.add(chart)
    return drawing


def build_excel_export(
    company: Dict[str, Any],
    filing: Dict[str, Any],
    history: pd.DataFrame,
    peers: pd.DataFrame,
    metric_sources: pd.DataFrame,
    sections: Dict[str, str],
    period_summary: str = "",
    disclaimer_lines: List[str] | None = None,
) -> bytes:
    buffer = BytesIO()
    with pd.ExcelWriter(buffer, engine="openpyxl") as writer:
        summary = pd.DataFrame(
            [
                ["View", "10K Summary"],
                ["Company", company.get("name")],
                ["Ticker", company.get("ticker")],
                ["Profile", company.get("profile")],
                ["Latest Filing Date", filing.get("filing_date")],
                ["Form", filing.get("form")],
                ["Filing URL", filing.get("filing_url")],
                ["Period Summary", period_summary],
            ],
            columns=["Field", "Value"],
        )
        summary.to_excel(writer, sheet_name="Summary", index=False)
        history.to_excel(writer, sheet_name="History", index=False)
        metric_sources.to_excel(writer, sheet_name="Metric Sources", index=False)
        if not peers.empty:
            peers.to_excel(writer, sheet_name="Peers", index=False)
        section_rows = [{"Section": key, "Text": value} for key, value in sections.items() if value]
        pd.DataFrame(section_rows).to_excel(writer, sheet_name="Sections", index=False)
        disclaimer_rows = [{"Item": idx + 1, "Disclaimer": line} for idx, line in enumerate(disclaimer_lines or [])]
        if disclaimer_rows:
            pd.DataFrame(disclaimer_rows).to_excel(writer, sheet_name="Disclaimer", index=False)
        _add_excel_charts(writer, history)
    return buffer.getvalue()


def build_pdf_export(
    company: Dict[str, Any],
    filing: Dict[str, Any],
    business_summary: str,
    period_summary: str,
    history: pd.DataFrame,
    kpi_rows: List[List[str]],
    focus_lines: List[str],
    peer_lines: List[str],
    reporting_lines: List[str],
    disclaimer_lines: List[str] | None = None,
) -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(A4),
        leftMargin=24,
        rightMargin=24,
        topMargin=22,
        bottomMargin=22,
    )
    styles = getSampleStyleSheet()
    title = styles["Title"]
    title.fontName = "Helvetica-Bold"
    title.fontSize = 18
    heading = styles["Heading2"]
    heading.fontName = "Helvetica-Bold"
    heading.fontSize = 11
    body = styles["BodyText"]
    body.fontName = "Helvetica"
    body.fontSize = 9
    body.leading = 12

    story: List[Any] = []
    story.append(Paragraph(f"10K Summary: {company.get('name')} ({company.get('ticker')})", title))
    story.append(Paragraph(period_summary, body))
    story.append(Spacer(1, 8))
    story.append(Paragraph("Business Overview", heading))
    story.append(Paragraph(business_summary, body))
    story.append(Spacer(1, 8))

    table = Table([["Metric", "Value", "Period"]] + kpi_rows, colWidths=[140, 110, 150])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#14213d")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE", (0, 0), (-1, -1), 8.5),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f6f8fb")]),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#d9e0ea")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]
        )
    )
    story.append(table)
    story.append(Spacer(1, 8))

    revenue_chart = _build_pdf_line_chart(
        history,
        "Revenue & Net Income Trend",
        ["Revenue", "Net Income"],
        ["#17365d", "#58718a"],
    )
    if revenue_chart is not None:
        story.append(revenue_chart)
        story.append(Spacer(1, 6))

    expense_chart = _build_pdf_line_chart(
        history,
        "Expenses Trend",
        ["Expenses"],
        ["#9a6700"],
    )
    if expense_chart is not None:
        story.append(expense_chart)
        story.append(Spacer(1, 6))

    balance_chart = _build_pdf_bar_chart(history, "Assets vs Liabilities", "Assets", "Liabilities")
    if balance_chart is not None:
        story.append(balance_chart)
        story.append(Spacer(1, 8))

    def bullets(title_text: str, lines: List[str]) -> None:
        story.append(Paragraph(title_text, heading))
        if lines:
            for line in lines[:4]:
                story.append(Paragraph(f"&bull; {line}", body))
        else:
            story.append(Paragraph("&bull; Not available.", body))
        story.append(Spacer(1, 6))

    bullets("Business Focus & Insights", focus_lines)
    bullets("Peer Positioning", peer_lines)
    bullets("Timeframe & Sources", reporting_lines)
    bullets("Disclaimer", (disclaimer_lines or [])[:4])

    doc.build(story)
    return buffer.getvalue()


def build_ppt_export(
    company: Dict[str, Any],
    business_summary: str,
    period_summary: str,
    history: pd.DataFrame,
    peers: pd.DataFrame,
    focus_lines: List[str],
    reporting_lines: List[str],
    disclaimer_lines: List[str] | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(8.2), Inches(0.6)).text_frame
    title.text = f"10K Summary | {company.get('name')} ({company.get('ticker')})"
    title.paragraphs[0].font.size = Pt(24)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.name = "Aptos"

    sub = slide.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(12.0), Inches(0.5)).text_frame
    sub.text = period_summary
    sub.paragraphs[0].font.size = Pt(10)
    sub.paragraphs[0].font.name = "Aptos"

    bus = slide.shapes.add_textbox(Inches(0.4), Inches(1.25), Inches(6.2), Inches(1.1)).text_frame
    bus.word_wrap = True
    bus.text = business_summary
    bus.paragraphs[0].font.size = Pt(11)
    bus.paragraphs[0].font.name = "Aptos"

    focus_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.25), Inches(5.8), Inches(2.0)).text_frame
    focus_box.text = "Business Focus & Insights"
    focus_box.paragraphs[0].font.size = Pt(14)
    focus_box.paragraphs[0].font.bold = True
    for line in focus_lines[:4]:
        p = focus_box.add_paragraph()
        p.text = line
        p.level = 1
        p.font.size = Pt(10)
        p.font.name = "Aptos"

    if not history.empty and {"Year", "Revenue", "Net Income"}.issubset(history.columns):
        chart_data = CategoryChartData()
        chart_data.categories = [str(int(year)) for year in history["Year"]]
        chart_data.add_series("Revenue ($B)", [float(value) / 1_000_000_000 for value in history["Revenue"]])
        chart_data.add_series("Net Income ($B)", [float(value) / 1_000_000_000 for value in history["Net Income"]])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.4),
            Inches(2.5),
            Inches(6.0),
            Inches(3.4),
            chart_data,
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Revenue & Net Income Trend"
        chart.value_axis.has_major_gridlines = True

    if not peers.empty:
        peer_chart_data = CategoryChartData()
        peer_chart_data.categories = list(peers["Ticker"].head(3))
        peer_chart_data.add_series("Revenue Growth %", [0 if pd.isna(v) else float(v) for v in peers["Revenue Growth %"].head(3)])
        peer_chart_data.add_series("Operating Margin %", [0 if pd.isna(v) else float(v) for v in peers["Operating Margin %"].head(3)])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(6.8),
            Inches(3.35),
            Inches(5.8),
            Inches(2.55),
            peer_chart_data,
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Peer Snapshot"

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title2 = slide2.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(12.2), Inches(0.5)).text_frame
    title2.text = "Reporting Basis & Notes"
    title2.paragraphs[0].font.size = Pt(22)
    title2.paragraphs[0].font.bold = True
    title2.paragraphs[0].font.name = "Aptos"

    box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.0), Inches(5.8)).text_frame
    box.text = "Timeframe & Sources"
    box.paragraphs[0].font.size = Pt(14)
    box.paragraphs[0].font.bold = True
    for line in reporting_lines[:4]:
        p = box.add_paragraph()
        p.text = line
        p.level = 1
        p.font.size = Pt(11)
        p.font.name = "Aptos"

    box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.0), Inches(2.5)).text_frame
    box2.text = "Source & Disclaimer"
    box2.paragraphs[0].font.size = Pt(14)
    box2.paragraphs[0].font.bold = True
    p = box2.add_paragraph()
    p.text = "Generated from the latest available Form 10-K plus external market data used for peer comparison."
    p.level = 1
    p.font.size = Pt(11)
    p.font.name = "Aptos"
    for line in (disclaimer_lines or [])[:3]:
        p = box2.add_paragraph()
        p.text = line
        p.level = 1
        p.font.size = Pt(10)
        p.font.name = "Aptos"

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_kpi_rows(history: pd.DataFrame) -> List[List[str]]:
    if history.empty:
        return []
    latest_row = history.iloc[-1]
    year = int(latest_row["Year"])
    rows = []
    for metric in ["Revenue", "Expenses", "Net Income", "Operating Cash Flow", "Free Cash Flow", "Diluted EPS", "Basic EPS"]:
        if metric not in history.columns:
            continue
        value = latest_row.get(metric)
        if pd.isna(value):
            continue
        if "EPS" in metric:
            display = f"{float(value):.2f}"
        elif "%" in metric:
            display = _fmt_percent(value)
        else:
            display = _fmt_currency(value)
        rows.append([metric, display, f"FY{year}"])
    return rows[:8]
